"""The business presentation is checked like code, because it is generated like code (US-38).

`scripts/build_presentation.py` turns the pipeline's artifacts into `docs/presentation.pptx`. The
deck is a course deliverable and a management-facing document, so the two ways it can go wrong both
have to be caught mechanically:

* **it stops matching the specification** — a slide dropped, the §53.2 order shuffled, the count
  outside the 10–12 the course brief allows, a "TODO" left on a slide;
* **it stops matching the numbers** — the far more dangerous failure, because a deck full of
  plausible stale figures looks exactly like a correct one. So the model-comparison table is
  compared cell by cell against `holdout_metrics_overall.csv`, and the run id printed on every
  slide is compared against `run_log.json` rather than merely being checked for existence: a deck
  rebuilt against tables from another run is precisely what that assertion catches.

The deck is rebuilt into a temporary directory rather than over `docs/`, so running the suite never
rewrites a committed deliverable. One test does check the committed file itself — it is the
artifact the course is graded on, and its absence must fail here rather than in a submission.
"""

from __future__ import annotations

import importlib.util
import json
import sys
import zipfile
from collections.abc import Iterator
from pathlib import Path
from types import ModuleType

import pytest
from pptx import Presentation

from pipeline import paths
from pipeline.config import MODEL_IDS

SCRIPT = paths.PROJECT_ROOT / "scripts" / "build_presentation.py"
COMMITTED_DECK = paths.DOCS_DIR / "presentation.pptx"


@pytest.fixture(scope="module")
def builder() -> ModuleType:
    """Import the generator by path — `scripts/` is not a package.

    The module is registered in ``sys.modules`` *before* it is executed, because ``@dataclass``
    resolves a field's annotation by looking its own module up there.
    """
    spec = importlib.util.spec_from_file_location("build_presentation", SCRIPT)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(scope="module")
def built(builder: ModuleType, tmp_path_factory: pytest.TempPathFactory) -> Iterator[Path]:
    """Build the deck from the current artifacts into a temporary directory."""
    output = tmp_path_factory.mktemp("deck")
    with pytest.MonkeyPatch.context() as patch:
        patch.setattr(builder, "PRESENTATION", output / "presentation.pptx")
        patch.setattr(builder, "PRESENTATION_NOTES", output / "presentation_notes.md")
        patch.setattr(builder, "DECK_FIGURES_DIR", output / "img" / "deck")
        yield builder.build()


@pytest.fixture(scope="module")
def deck(built: Path) -> Presentation:
    return Presentation(str(built))


@pytest.fixture(scope="module")
def notes(builder: ModuleType, built: Path) -> str:
    return builder.PRESENTATION_NOTES.read_text(encoding="utf-8")


@pytest.fixture(scope="module")
def run_log() -> dict:
    return json.loads(paths.RUN_LOG.read_text(encoding="utf-8"))


def _slide_text(slide) -> str:  # noqa: ANN001 - a pptx Slide
    """Every character a reader can see on one slide, tables included."""
    parts: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
        if getattr(shape, "has_table", False) and shape.has_table:
            parts.extend(cell.text for row in shape.table.rows for cell in row.cells)
    return "\n".join(parts)


def _title(slide) -> str:  # noqa: ANN001
    """The slide's title is the first shape the builder adds."""
    return slide.shapes[0].text_frame.text


def _pictures(slide) -> list:  # noqa: ANN001
    return [shape for shape in slide.shapes if shape.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE


# ---------------------------------------------------------------------------
# It runs, and it produces both deliverables
# ---------------------------------------------------------------------------


def test_the_generator_runs_on_the_current_artifacts(builder: ModuleType, built: Path) -> None:
    assert built.is_file(), "build() reported success but wrote no deck"
    assert built.stat().st_size > 0
    assert builder.PRESENTATION_NOTES.is_file(), "the speaker notes were not written"


def test_two_builds_of_the_same_run_are_byte_identical(
    builder: ModuleType, built: Path, tmp_path: Path
) -> None:
    """§40's property, applied to a committed generated file.

    A `.pptx` is a zip, and python-pptx stamps its members with the clock, so without the archive
    rewrite two identical decks differ in their bytes and `git status` reports a change that is not
    one. This is the test that keeps that rewrite in place.
    """
    with pytest.MonkeyPatch.context() as patch:
        rebuilt = tmp_path / "presentation.pptx"
        patch.setattr(builder, "PRESENTATION", rebuilt)
        patch.setattr(builder, "PRESENTATION_NOTES", tmp_path / "presentation_notes.md")
        patch.setattr(builder, "DECK_FIGURES_DIR", tmp_path / "img" / "deck")
        builder.build()
    assert rebuilt.read_bytes() == built.read_bytes()


def test_the_deck_is_a_valid_readable_archive(built: Path) -> None:
    """PowerPoint reads `[Content_Types].xml` first; the archive rewrite must not reorder it."""
    with zipfile.ZipFile(built) as archive:
        assert archive.namelist()[0] == "[Content_Types].xml"
        assert archive.testzip() is None


def test_the_committed_deliverable_exists_and_is_the_right_size(builder: ModuleType) -> None:
    """The deck in `docs/` is what gets submitted; criterion 22 of §49 reads it from there."""
    assert COMMITTED_DECK.is_file(), f"{COMMITTED_DECK} is missing — generate it and commit it"
    slides = len(Presentation(str(COMMITTED_DECK)).slides)
    assert builder.SLIDE_COUNT_MIN <= slides <= builder.SLIDE_COUNT_MAX


# ---------------------------------------------------------------------------
# It matches the §53.2 outline
# ---------------------------------------------------------------------------


def test_slide_count_is_within_the_course_brief_bounds(
    builder: ModuleType, deck: Presentation
) -> None:
    assert builder.SLIDE_COUNT_MIN <= len(deck.slides) <= builder.SLIDE_COUNT_MAX


def test_slide_titles_follow_the_53_2_outline_in_order(
    builder: ModuleType, deck: Presentation
) -> None:
    assert [_title(slide) for slide in deck.slides] == list(builder.SLIDE_TITLES)


def test_no_slide_carries_a_placeholder(deck: Presentation) -> None:
    for number, slide in enumerate(deck.slides, start=1):
        assert "TODO" not in _slide_text(slide), f"slide {number} still contains a placeholder"


# ---------------------------------------------------------------------------
# It matches the numbers
# ---------------------------------------------------------------------------


def test_every_slide_names_the_audited_run_and_the_dataset_licence(
    deck: Presentation, run_log: dict
) -> None:
    """Not "a run id appears" but "*this* run's id appears" — the stale-deck check."""
    for number, slide in enumerate(deck.slides, start=1):
        text = _slide_text(slide)
        assert run_log["run_id"] in text, f"slide {number} does not name run {run_log['run_id']}"
        assert "CC BY 4.0" in text, f"slide {number} is missing the dataset licence"


def test_slide_7_reproduces_holdout_metrics_overall(
    builder: ModuleType, deck: Presentation
) -> None:
    """Every candidate, with the wMAPE and Bias the evaluation table holds — §23 wants both."""
    import pandas as pd

    metrics = pd.read_csv(paths.EVAL_TABLES_DIR / "holdout_metrics_overall.csv")
    text = _slide_text(list(deck.slides)[6])
    for model_id in MODEL_IDS:
        row = metrics.loc[metrics["model"] == model_id].iloc[0]
        assert model_id in text, f"slide 7 does not list candidate {model_id}"
        assert builder._pct(row["wmape"]) in text, f"slide 7 has the wrong wMAPE for {model_id}"
        assert builder._pct(row["bias"]) in text, f"slide 7 has the wrong Bias for {model_id}"


def test_slide_8_states_the_champion_and_the_gate_outcome(deck: Presentation) -> None:
    decision = json.loads((paths.CHAMPION_DECISION).read_text(encoding="utf-8"))
    text = _slide_text(list(deck.slides)[7])
    assert decision["champion"] in text
    assert "Gate 1" in text and "Gate 2" in text
    # The verdict must be stated either way round: an ML win and a baseline win are both legitimate.
    verdict = "meaningful improvement" if decision["meaningful_improvement"] else "competitive"
    assert verdict in text.lower()


def test_slide_9_compares_champion_and_baseline_on_the_three_kpis(
    builder: ModuleType, deck: Presentation, built: Path
) -> None:
    numbers = builder.Numbers.load(builder.RunProvenance.load())
    text = _slide_text(list(deck.slides)[8])
    for model_id in (numbers.champion_id, numbers.baseline_id):
        assert model_id in text
        for column, render in (
            ("fill_rate", builder._pct),
            ("stockout_units", builder._units),
            ("excess_units", builder._units),
        ):
            value = numbers.kpi(model_id, builder._POLICY_WITH_SAFETY_STOCK, column)
            assert render(value) in text, f"slide 9 is missing {column} for {model_id}"


def test_slide_9_carries_the_configured_service_level_disclaimer(
    builder: ModuleType, deck: Presentation
) -> None:
    """PRD §25: z does not guarantee a fill rate, and the deck must not imply otherwise."""
    numbers = builder.Numbers.load(builder.RunProvenance.load())
    assert numbers.disclaimer in _slide_text(list(deck.slides)[8])


def test_slide_10_embeds_the_committed_screenshots(builder: ModuleType, deck: Presentation) -> None:
    pictures = _pictures(list(deck.slides)[9])
    assert len(pictures) >= 3, "slide 10 must show at least three screens of the product"
    committed = {path.read_bytes() for path in builder.SCREENS_DIR.glob("*.png")}
    for picture in pictures:
        assert picture.image.blob in committed, "a slide-10 image is not a committed screenshot"


# ---------------------------------------------------------------------------
# The speaker notes
# ---------------------------------------------------------------------------


def test_the_notes_have_one_section_per_slide(deck: Presentation, notes: str) -> None:
    assert notes.count("\n## Slide ") == len(deck.slides)


def test_the_notes_never_promise_an_order_amount(notes: str) -> None:
    """PRD §7: the output is a Recommended Target Inventory. The AC greps for this exact string."""
    assert "Order Quantity" not in notes


def test_every_speaker_note_is_within_the_word_budget(builder: ModuleType, notes: str) -> None:
    sections = notes.split("\n## Slide ")[1:]
    assert sections, "the notes file has no slide sections"
    for section in sections:
        heading, _, body = section.partition("\n")
        assert len(body.split()) <= builder.MAX_NOTE_WORDS, f"note for slide {heading} is too long"


def test_the_notes_name_the_audited_run(notes: str, run_log: dict) -> None:
    assert run_log["run_id"] in notes
    assert run_log["status"] in notes


# ---------------------------------------------------------------------------
# The provenance rules from the issue's §8 interface corrections
# ---------------------------------------------------------------------------


def test_a_run_that_did_not_succeed_is_refused(
    builder: ModuleType, tmp_path: Path, monkeypatch: pytest.MonkeyPatch, run_log: dict
) -> None:
    """A deck built from a failed run would show the previous run's tables under this run's id."""
    broken = tmp_path / "run_log.json"
    broken.write_text(json.dumps({**run_log, "status": "failed"}), encoding="utf-8")
    monkeypatch.setattr(paths, "RUN_LOG", broken)
    with pytest.raises(SystemExit, match="not 'success'"):
        builder.build()


def test_a_preliminary_build_stamps_every_slide(builder: ModuleType) -> None:
    """`--allow-preliminary` is only safe because the stamp is impossible to miss."""
    stamp = "PRELIMINARY — run test, status=running"
    deck = builder.Deck(footnote="footnote", stamp=stamp)
    slide = deck.slide("Title")
    assert stamp in _slide_text(slide)


def test_undeclared_inputs_are_disclosed_not_hidden(builder: ModuleType, notes: str) -> None:
    """A table the audited run did not write is named, so nobody presents it as this run's."""
    run = builder.RunProvenance.load()
    builder.Numbers.load(run)
    if run.has_provenance_gap:
        assert "Provenance warning" in notes
        for undeclared in run.undeclared:
            assert undeclared in notes
    else:
        assert "Provenance warning" not in notes


def test_the_baseline_and_the_candidates_come_from_configuration(builder: ModuleType) -> None:
    """§20 gate 4 compares against the configured main baseline; the deck must not assume B2."""
    numbers = builder.Numbers.load(builder.RunProvenance.load())
    from pipeline.config import load_model_config

    assert numbers.baseline_id == load_model_config().main_baseline_id
    assert set(numbers.overall["model"]) >= set(MODEL_IDS)
