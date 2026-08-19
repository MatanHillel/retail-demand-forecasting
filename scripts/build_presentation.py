"""Build the business presentation from the pipeline's own artifacts (US-38, PRD §53.2).

The course brief asks for a 10–12 slide business presentation. PRD §53.2 fixes its outline, and the
project's §14 convention forbids a final number from being typed by a human: every figure a reader
sees must have been computed by the pipeline and written to a table. A deck maintained by hand
breaks that rule the first time a model is retrained, so the deck is *generated*. This script reads
`artifacts/`, `config/` and `docs/PRD.md`, and writes:

* ``docs/presentation.pptx`` — twelve slides in the §53.2 order;
* ``docs/presentation_notes.md`` — the speaker notes, one section per slide;
* ``docs/img/deck/*.png`` — the four figures the deck needs that the pipeline does not already draw.

Run it after a pipeline run::

    python scripts/build_presentation.py
    python scripts/build_presentation.py --pdf              # also export a PDF, if LibreOffice is
    python scripts/build_presentation.py --allow-preliminary  # build from a non-successful run

Three rules govern where the numbers come from (they are the §8 interface corrections on AI-43):

1. **The run id, the timestamp and the status come from ``run_log.json``, and nothing is
   recomputed.** ``status`` has three values — ``running``, ``success``, ``failed`` — and a deck
   built from a run that did not succeed would present the *previous* run's tables under the
   current run's id. This script therefore refuses to build unless the run succeeded, unless
   ``--allow-preliminary`` is passed, which stamps "PRELIMINARY" on every slide instead.
2. **A table is only attributed to the audited run when the run log says so.** ``promote()`` is
   refused once a run fails, so the files in their final locations can belong to an older run than
   the one named in the log. Every input is checked against ``run_log.json → artifacts`` and against
   the table's own recorded run id where it has one; anything that does not match is listed in a
   provenance note on slide 1 and in the notes file. The deck never silently presents a stale table
   as this run's result.
3. **Configuration comes from the run's snapshot, not from the current YAML.** ``k`` on slide 5 and
   ``z`` on slide 9 are labelled "as of run X", and ``config/*.yaml`` may have moved since; so they
   are read from ``run_log.json → config_snapshot``, falling back to the typed loaders
   (:func:`~pipeline.config.load_model_config`, :func:`~pipeline.config.load_inventory_policy`) only
   when the snapshot is absent.

The script opens **no** :class:`~pipeline.run_context.RunContext`: ``start()`` would allocate a new
run id — which slide 1 would then display instead of the pipeline's — and ``finish()`` would
overwrite the very ``run_log.json`` the deck is quoting. It writes under ``docs/`` only, so it needs
no staging path either (§39 applies to artifacts, and these are deliverables).

The only numbers written in this file are layout: inches, points and colours.
"""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
import zipfile
from collections.abc import Iterable, Sequence
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path

_REPO_ROOT = Path(__file__).resolve().parents[1]
if str(_REPO_ROOT / "src") not in sys.path:  # the script may run without an editable install
    sys.path.insert(0, str(_REPO_ROOT / "src"))

import matplotlib.pyplot as plt  # noqa: E402  (must follow the sys.path bootstrap)
import pandas as pd  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Emu, Inches, Pt  # noqa: E402

from pipeline import paths  # noqa: E402
from pipeline.config import MODEL_IDS, load_inventory_policy, load_model_config  # noqa: E402
from pipeline.eda import style  # noqa: E402

# ---------------------------------------------------------------------------
# Output locations
# ---------------------------------------------------------------------------

PRESENTATION: Path = paths.DOCS_DIR / "presentation.pptx"
PRESENTATION_NOTES: Path = paths.DOCS_DIR / "presentation_notes.md"
PRESENTATION_PDF: Path = paths.DOCS_DIR / "presentation.pdf"
SCREENS_DIR: Path = paths.DOCS_DIR / "img" / "screens"
DECK_FIGURES_DIR: Path = paths.DOCS_DIR / "img" / "deck"
PRD: Path = paths.DOCS_DIR / "PRD.md"
FLOW_DOC: Path = paths.DOCS_DIR / "flow.md"

#: PRD §53.2 allows 10–12 slides; the course brief says the same. Asserted before saving.
SLIDE_COUNT_MIN: int = 10
SLIDE_COUNT_MAX: int = 12

#: A speaker note is a prompt, not a script. §2 of the issue caps them.
MAX_NOTE_WORDS: int = 120

# ---------------------------------------------------------------------------
# Layout — the only numbers this file is allowed to contain
# ---------------------------------------------------------------------------

SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.62)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN
TITLE_TOP = Inches(0.34)
TITLE_HEIGHT = Inches(0.62)
KICKER_TOP = Inches(0.96)
KICKER_HEIGHT = Inches(0.36)
BODY_TOP = Inches(1.44)
BODY_HEIGHT = Inches(5.24)
FOOTNOTE_TOP = Inches(6.86)
FOOTNOTE_HEIGHT = Inches(0.34)
RULE_TOP = Inches(1.28)
RULE_HEIGHT = Pt(1.5)
GUTTER = Inches(0.34)
TABLE_HEIGHT = Inches(3.6)

TITLE_SIZE = Pt(30)
KICKER_SIZE = Pt(14)
BODY_SIZE = Pt(15)
SMALL_SIZE = Pt(12)
TABLE_SIZE = Pt(12)
FOOTNOTE_SIZE = Pt(9)
STAMP_SIZE = Pt(11)

INK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5A, 0x5A, 0x5A)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
BAND = RGBColor(0xF2, 0xF4, 0xF6)
WARNING = RGBColor(0xD5, 0x5E, 0x00)  # Okabe–Ito vermillion, the palette's alert colour


def _rgb(hex_colour: str) -> RGBColor:
    """``"#0072B2"`` -> an ``RGBColor``. The palette is defined once, in `pipeline.eda.style`."""
    return RGBColor.from_string(hex_colour.lstrip("#").upper())


ACCENT = _rgb(style.PALETTE[0])

#: A share (0–1) is shown to a reader as a percentage; a difference of two shares is shown in
#: "points". Both are unit conversions applied at format time, never to a stored value.
PERCENT_SCALE = 100.0


def _pct(share: float, decimals: int = 1) -> str:
    """``0.557127`` -> ``"55.7 %"``."""
    return f"{share * PERCENT_SCALE:.{decimals}f} %"


def _points(share_difference: float, decimals: int = 2) -> str:
    """``0.021454`` -> ``"+2.15"`` — a difference between two shares, in percentage points."""
    return f"{share_difference * PERCENT_SCALE:+.{decimals}f}"


def _units(value: float) -> str:
    """``398498.0`` -> ``"398,498"`` — units are counts and are never shown with decimals."""
    return f"{value:,.0f}"


# ---------------------------------------------------------------------------
# The audited run
# ---------------------------------------------------------------------------


@dataclass
class RunProvenance:
    """What ``run_log.json`` says about the run whose numbers this deck presents.

    ``artifacts`` is the run's own ``{key: repo-relative path}`` map. A file that is *not* in it was
    not written by this run: because :meth:`~pipeline.run_context.RunContext.promote` is refused
    after a failure, the file in the final location may be an older run's. :meth:`attribute` records
    that judgement per input so the deck can disclose it instead of hiding it.
    """

    run_id: str
    status: str
    started_at: str
    finished_at: str
    mode: str
    champion: dict | None
    config_snapshot: dict
    artifacts: dict[str, str]
    #: Inputs read by the deck that the audited run did not declare, in first-read order.
    undeclared: list[str] = field(default_factory=list)
    #: ``{input: run id}`` for inputs that carry a run id of their own that is not this run's.
    foreign_run_ids: dict[str, str] = field(default_factory=dict)

    @classmethod
    def load(cls) -> RunProvenance:
        if not paths.RUN_LOG.is_file():
            raise SystemExit(
                f"{_rel(paths.RUN_LOG)} does not exist — the deck quotes a run, so there has to "
                "be one. Run `python -m pipeline --no-llm` first."
            )
        log = json.loads(paths.RUN_LOG.read_text(encoding="utf-8"))
        return cls(
            run_id=str(log.get("run_id", "")),
            status=str(log.get("status", "")),
            started_at=str(log.get("started_at", "")),
            finished_at=str(log.get("finished_at", "")),
            mode=str(log.get("mode", "")),
            champion=log.get("champion"),
            config_snapshot=log.get("config_snapshot") or {},
            artifacts={str(k): str(v) for k, v in (log.get("artifacts") or {}).items()},
        )

    @property
    def succeeded(self) -> bool:
        return self.status == "success"

    @property
    def date(self) -> str:
        """The date of the numbers: the run's finish, falling back to its start."""
        stamp = self.finished_at or self.started_at
        return stamp.split("T", maxsplit=1)[0] if stamp else "unknown date"

    @property
    def archive_timestamp(self) -> tuple[int, int, int, int, int, int]:
        """The run's finish time as a zip timestamp — see :func:`_freeze_archive_timestamps`."""
        stamp = self.finished_at or self.started_at
        try:
            moment = datetime.fromisoformat(stamp)
        except ValueError:
            return _ZIP_EPOCH
        return (moment.year, moment.month, moment.day, moment.hour, moment.minute, moment.second)

    def attribute(self, path: Path, own_run_id: str | None = None) -> Path:
        """Record where ``path`` came from, and return it unchanged so reads can wrap the call.

        ``own_run_id`` is the run id the file records about itself, when it has one (the tables
        that carry a ``run_id`` column or key). It is the stronger evidence: a file can be listed
        in the run log and still be an older copy only if something went badly wrong, whereas a
        recorded id that differs from the audited run's is proof of exactly that.
        """
        relative = _rel(path)
        if relative not in self.artifacts.values() and relative not in self.undeclared:
            self.undeclared.append(relative)
        if own_run_id and own_run_id != self.run_id:
            self.foreign_run_ids[relative] = own_run_id
        return path

    @property
    def has_provenance_gap(self) -> bool:
        return bool(self.undeclared or self.foreign_run_ids)

    def provenance_sentence(self) -> str:
        """One sentence naming the gap, for slide 1 and the notes file. Empty when there is none.

        The caller supplies its own label, so the sentence reads correctly after both "Provenance:"
        on the slide and "Provenance warning." in the notes.
        """
        if not self.has_provenance_gap:
            return ""
        parts = [
            f"{len(self.undeclared)} input table(s) are not listed in this run's artifact map, so "
            "they were carried forward from an earlier run"
        ]
        if self.foreign_run_ids:
            ids = ", ".join(sorted(set(self.foreign_run_ids.values())))
            parts.append(f"and record run id(s) {ids} of their own")
        return " ".join(parts) + "."

    def snapshot(self, config_name: str) -> dict:
        """One config file as it was at run time; ``{}`` when the run recorded no snapshot."""
        section = self.config_snapshot.get(config_name)
        return section if isinstance(section, dict) else {}


def _rel(path: Path) -> str:
    """Repo-relative, forward-slashed — the form ``run_log.json`` records paths in."""
    try:
        return path.resolve().relative_to(paths.PROJECT_ROOT).as_posix()
    except ValueError:  # pragma: no cover - a path outside the repository
        return path.as_posix()


# ---------------------------------------------------------------------------
# Reading the documents the deck quotes prose from
# ---------------------------------------------------------------------------


def _prd_section(number: int) -> str:
    """The body of one numbered PRD section, so the deck's prose cannot drift from the spec."""
    text = PRD.read_text(encoding="utf-8")
    match = re.search(
        rf"^# {number}\. [^\n]*\n(.*?)(?=^# \d+\.|^# Appendix)", text, re.MULTILINE | re.DOTALL
    )
    if match is None:  # pragma: no cover - only if the PRD is restructured
        raise SystemExit(f"PRD §{number} not found in {_rel(PRD)}")
    return match.group(1).strip()


def _prd_quote(number: int) -> str:
    """The first block-quoted sentence of a PRD section — §1's question, §4's vision, §51's MVP."""
    for line in _prd_section(number).splitlines():
        if line.startswith(">"):
            return line.lstrip("> ").replace("**", "").strip()
    raise SystemExit(f"PRD §{number} has no block quote")  # pragma: no cover


def _flow_steps() -> list[str]:
    """The ten Flow step names, read out of `docs/flow.md`'s step diagram (§37)."""
    diagram = FLOW_DOC.read_text(encoding="utf-8")
    steps = re.findall(r"^\s*→\s+(\d+)\s+(\w+)", diagram, re.MULTILINE)
    return [f"{number}. {name.replace('_', ' ')}" for number, name in steps]


def _crew_agents(crew: str) -> list[str]:
    """The three agent roles of one crew, read from its own `agents.yaml` (§35–§36)."""
    import yaml

    config = _REPO_ROOT / "src" / "crews" / crew / "config" / "agents.yaml"
    agents = yaml.safe_load(config.read_text(encoding="utf-8"))
    return [" ".join(str(spec["role"]).split()) for spec in agents.values()]


def _team() -> list[tuple[str, str]]:
    """Team members and the §54 roles they own, names read from `pyproject.toml`."""
    import tomllib

    pyproject = tomllib.loads((_REPO_ROOT / "pyproject.toml").read_text(encoding="utf-8"))
    names = [author["name"] for author in pyproject["project"]["authors"]]
    # §54 lists five roles for up to five students; a three-person team doubles up. The mapping
    # follows the ownership split recorded in CLAUDE.md §8.
    roles = {
        "Daniel B. Fogel": "App lead · PM / QA / docs — platform, Streamlit, CI, presentation",
        "Matan Hillel": "Data & EDA lead · Flow lead — cleaning, panel, EDA, contract, Flow",
        "Dor Hll": "ML lead — features, models, evaluation, sigma, inventory policy",
    }
    return [(name, roles.get(name, "Contributor")) for name in names]


def _insight_sentences(tables: Sequence[str], max_chars: int) -> list[str]:
    """The numbered insight backed by each named EDA table — slide 4's three "so what" sentences.

    `insights.md` is written by the pipeline, and every number in it is checked against the tables
    under `eda_tables/` before it is accepted, in both `--no-llm` and LLM mode (§38). Quoting it,
    rather than composing a sentence here from the same tables, is what puts slide 4 inside that
    guarantee. An insight whose table is no longer cited is an error, not a blank bullet.
    """
    text = paths.INSIGHTS.read_text(encoding="utf-8")
    found: list[str] = []
    for table in tables:
        match = re.search(rf"^\d+\.\s+(.*?)\s*\(E\d+, table `{re.escape(table)}`\)\s*$", text, re.M)
        if match is None:
            raise SystemExit(
                f"insights.md has no insight citing table '{table}' — slide 4 quotes the "
                "pipeline's own insights and must not invent one"
            )
        found.append(_shorten(match.group(1).replace("**", "").strip(), max_chars))
    return found


def _markdown_bullets(body: str) -> list[str]:
    """Markdown bullets as whole sentences, one string each.

    `model_card.md` is written for reading, so its bullets wrap over several lines and some carry
    nested sub-bullets. Taking only the lines that begin with a dash would truncate every bullet
    mid-sentence, so continuation lines are joined back on and a nested bullet is folded into its
    parent.
    """
    bullets: list[str] = []
    for raw in body.splitlines():
        stripped = raw.strip()
        if not stripped:
            continue
        is_bullet = stripped.startswith(("- ", "* "))
        indented = raw[:1].isspace()
        # Emphasis and code ticks are dropped as the line is read, not at the end: the fold below
        # decides on punctuation, and "**Residual risks:**" does not end in a colon until it has.
        content = re.sub(r"[*`]", "", stripped[2:] if is_bullet else stripped).strip()
        if is_bullet and not indented:
            bullets.append(content)
        elif not bullets:
            continue
        elif is_bullet:  # a nested sub-bullet belongs to the bullet above it
            parent = bullets[-1]
            bullets[-1] = f"{parent}{' ' if parent.endswith(':') else '; '}{content}"
        else:  # a wrapped continuation line
            bullets[-1] = f"{bullets[-1]} {content}"
    return bullets


#: A sentence-boundary trim is only used when it keeps at least this share of the allowance;
#: otherwise a bullet whose first sentence is three words long would lose everything that matters.
_MIN_KEPT_SHARE = 0.6


def _shorten(text: str, max_chars: int) -> str:
    """Trim to a whole sentence or clause — a slide bullet that stops mid-word reads as a bug."""
    if len(text) <= max_chars:
        return text
    head = text[:max_chars]
    stop = max(head.rfind(". "), head.rfind("; "))
    if stop > max_chars * _MIN_KEPT_SHARE:
        return head[:stop].strip().rstrip(",;:") + "."
    return f"{head[: head.rfind(' ')].rstrip(' ,;:')}..."


def _model_card_bullets(section_title: str, topics: Sequence[str], max_chars: int) -> list[str]:
    """Selected bullets of one `model_card.md` section — slide 11's limitations and ethics.

    ``topics`` names the opening words of the bullets the slide has room for, in the order §47/§48
    lists them; a topic the model card no longer carries is a hard error rather than a silent gap,
    because slide 11 is the honesty slide and a missing limitation is exactly what must not vanish.
    """
    text = paths.MODEL_CARD.read_text(encoding="utf-8")
    match = re.search(
        rf"^## \d+\. {re.escape(section_title)}\n(.*?)(?=^## |\Z)", text, re.MULTILINE | re.DOTALL
    )
    if match is None:  # pragma: no cover - only if US-25's template changes
        raise SystemExit(f"model_card.md has no section '{section_title}'")
    bullets = _markdown_bullets(match.group(1))
    selected: list[str] = []
    for topic in topics:
        lead = topic.lower()
        found = next((b for b in bullets if b.lower().startswith(lead)), None)
        if found is None:
            raise SystemExit(
                f"model_card.md section '{section_title}' has no bullet starting '{topic}' — "
                "slide 11 quotes the model card and must not drop a limitation silently"
            )
        selected.append(_shorten(found, max_chars))
    return selected


def _roadmap_items(limit: int) -> list[str]:
    """PRD §50's version-2 roadmap, split into its semicolon-separated items."""
    body = " ".join(_prd_section(50).split())
    return [item.strip(" .").capitalize() for item in body.split(";")][:limit]


# ---------------------------------------------------------------------------
# Reading the tables
# ---------------------------------------------------------------------------


@dataclass
class Numbers:
    """Every table the deck needs, read once, with its provenance already recorded."""

    run: RunProvenance
    quality: dict
    waterfall: pd.DataFrame
    overall: pd.DataFrame
    by_abc: pd.DataFrame
    improvement: pd.DataFrame
    champion: dict
    kpis: pd.DataFrame
    excess: pd.DataFrame
    plan: pd.DataFrame

    @classmethod
    def load(cls, run: RunProvenance) -> Numbers:
        quality_path = run.attribute(paths.DATA_QUALITY_FINDINGS)
        quality = json.loads(quality_path.read_text(encoding="utf-8"))
        run.attribute(quality_path, str(quality.get("run_id", "")))

        champion_path = run.attribute(paths.CHAMPION_DECISION)
        champion = json.loads(champion_path.read_text(encoding="utf-8"))
        run.attribute(champion_path, str(champion.get("run_id", "")))

        plan = pd.read_csv(run.attribute(paths.INVENTORY_PLAN), dtype={"stock_code": str})
        if "run_id" in plan.columns and not plan["run_id"].empty:
            run.attribute(paths.INVENTORY_PLAN, str(plan["run_id"].iloc[0]))

        # Slide 4 quotes `insights.md` and slide 11 quotes `model_card.md`; they are inputs like
        # any table, so they are attributed here even though they are read where they are used.
        run.attribute(paths.INSIGHTS)
        run.attribute(paths.MODEL_CARD)

        eda = paths.EDA_TABLES_DIR
        evaluation = paths.EVAL_TABLES_DIR
        return cls(
            run=run,
            quality=quality,
            waterfall=pd.read_csv(run.attribute(eda / "E01_cleaning_waterfall.csv")),
            overall=pd.read_csv(run.attribute(evaluation / "holdout_metrics_overall.csv")),
            by_abc=pd.read_csv(run.attribute(evaluation / "holdout_metrics_by_abc.csv")),
            improvement=pd.read_csv(run.attribute(evaluation / "improvement_vs_b2.csv")),
            champion=champion,
            kpis=pd.read_csv(run.attribute(paths.INVENTORY_KPIS)),
            excess=pd.read_csv(run.attribute(paths.EXCESS_CONCENTRATION)),
            plan=plan,
        )

    # -- derived views the slides ask for --------------------------------

    @property
    def champion_id(self) -> str:
        return str(self.champion["champion"])

    @property
    def baseline_id(self) -> str:
        """The main baseline (§20 gate 4 compares against it) — config, never a literal."""
        snapshot_models = self.run.snapshot("model_config").get("models", {})
        for model_id, spec in snapshot_models.items():
            if spec.get("main_baseline"):
                return str(model_id)
        return load_model_config().main_baseline_id

    @property
    def z(self) -> float:
        """The service-level multiplier as it stood at run time."""
        snapshot = self.run.snapshot("inventory_policy")
        return float(snapshot.get("z", load_inventory_policy().z))

    @property
    def z_options(self) -> list[float]:
        snapshot = self.run.snapshot("inventory_policy")
        return [float(z) for z in snapshot.get("z_options", load_inventory_policy().z_options)]

    @property
    def disclaimer(self) -> str:
        """The §25 sentence, from configuration rather than retyped into the template."""
        snapshot = self.run.snapshot("inventory_policy")
        return str(snapshot.get("disclaimer", load_inventory_policy().disclaimer))

    @property
    def active_k(self) -> int:
        snapshot = self.run.snapshot("model_config")
        active_rule = snapshot.get("active_rule", {})
        return int(active_rule.get("k", load_model_config().active_rule.k))

    @property
    def split(self):  # noqa: ANN201 - a pydantic model from pipeline.config
        return load_model_config().split

    def metric(self, model_id: str, column: str) -> float:
        row = self.overall.loc[self.overall["model"] == model_id]
        return float(row[column].iloc[0])

    def kpi(self, model_id: str, policy: str, column: str) -> float:
        """One overall back-test KPI, at the configured z."""
        rows = self.kpis.loc[
            (self.kpis["model"] == model_id)
            & (self.kpis["policy"] == policy)
            & (self.kpis["scope"] == "overall")
            & (self.kpis["z"].round(_Z_MATCH_DECIMALS) == round(self.z, _Z_MATCH_DECIMALS))
        ]
        return float(rows[column].iloc[0])

    def top_a_product(self) -> pd.Series:
        """The largest forecast among active class-A products — slide 2's worked example."""
        active = self.plan.loc[
            (self.plan["status"] == _FORECAST_STATUS) & (self.plan["abc_class"] == _CLASS_A)
        ]
        return active.sort_values("forecast", ascending=False).iloc[0]


#: Floating point z values are compared after rounding; ``1.645`` reads back exactly, but a policy
#: file could carry more decimals than the KPI table's ``%.6f`` format preserves.
_Z_MATCH_DECIMALS = 6
#: The `inventory_plan.csv` status that marks an active product (US-23's vocabulary).
_FORECAST_STATUS = "Forecast"
_CLASS_A = "A"
#: The two inventory policies compared in the back-test (§29).
_POLICY_FORECAST_ONLY = "forecast_only"
_POLICY_WITH_SAFETY_STOCK = "forecast_plus_ss"


# ---------------------------------------------------------------------------
# The four figures the deck draws for itself
# ---------------------------------------------------------------------------


def _save_deck_figure(figure, name: str) -> Path:
    """Save one deck figure under `docs/img/deck/` at the project's figure resolution."""
    DECK_FIGURES_DIR.mkdir(parents=True, exist_ok=True)
    path = DECK_FIGURES_DIR / f"{name}.png"
    figure.savefig(path, dpi=style.FIGURE_DPI, bbox_inches="tight", facecolor="white")
    plt.close(figure)
    return path


def draw_flow_diagram(steps: Sequence[str]) -> Path:
    """Slide 6's architecture picture: the ten Flow steps and the four validation routers (§37)."""
    style.apply_style()
    figure, axes = plt.subplots(figsize=(11.0, 4.6))
    axes.set_xlim(0, 1)
    axes.set_ylim(0, 1)
    axes.axis("off")

    rows = 5  # two columns of five steps
    box_width, box_height = 0.42, 0.13
    x_gap, y_gap = 0.06, 0.075
    router_after = {1, 3, 5, 9}  # the four @router checkpoints (§37)

    for index, label in enumerate(steps):
        column, row = divmod(index, rows)
        left = column * (box_width + x_gap)
        bottom = 1 - box_height - row * (box_height + y_gap)
        number = index + 1
        is_gate = number in router_after
        axes.add_patch(
            plt.Rectangle(
                (left, bottom),
                box_width,
                box_height,
                facecolor=style.PALETTE[0] if is_gate else "#F2F4F6",
                edgecolor=style.PALETTE[0],
                linewidth=1.2,
                alpha=0.95 if is_gate else 1.0,
            )
        )
        axes.text(
            left + box_width / 2,
            bottom + box_height / 2,
            label,
            ha="center",
            va="center",
            fontsize=11,
            color="white" if is_gate else "#1A1A1A",
            fontweight="bold" if is_gate else "normal",
        )
        if row < rows - 1:
            axes.annotate(
                "",
                xy=(left + box_width / 2, bottom - y_gap * 0.1),
                xytext=(left + box_width / 2, bottom - y_gap * 0.9),
                arrowprops={"arrowstyle": "<-", "color": "#5A5A5A", "linewidth": 1.1},
            )
    # The wrap arrow: the left column's last step continues into the right column's first.
    last_left_bottom = 1 - box_height - (rows - 1) * (box_height + y_gap)
    last_left_middle = last_left_bottom + box_height / 2
    axes.annotate(
        "",
        xy=(box_width + x_gap, 1 - box_height / 2),
        xytext=(box_width, last_left_middle),
        arrowprops={
            "arrowstyle": "->",
            "color": "#5A5A5A",
            "linewidth": 1.1,
            "connectionstyle": "arc3,rad=-0.45",
        },
    )
    axes.text(
        0.5,
        -0.04,
        "Filled steps are followed by a @router: it returns 'continue' or 'fail', and a failed run "
        "writes validation_report.json and promotes nothing (PRD §39).",
        ha="center",
        va="top",
        fontsize=10,
        color="#5A5A5A",
        transform=axes.transAxes,
    )
    figure.suptitle("CrewAI Flow — ten steps, four validation routers", fontsize=14, y=1.02)
    return _save_deck_figure(figure, "P01_flow_architecture")


def draw_target_inventory_example(numbers: Numbers) -> Path:
    """Slide 2's worked example: forecast, safety stock and target inventory for a real product."""
    product = numbers.top_a_product()
    forecast = float(product["forecast"])
    safety_stock = float(product["safety_stock"])
    target = float(product["target_inventory"])

    style.apply_style()
    figure, axes = plt.subplots(figsize=(9.0, 2.8))
    bar_height = 0.5
    label = "Recommended\nTarget Inventory"
    axes.barh([label], [forecast], height=bar_height, color=style.PALETTE[0])
    axes.barh([label], [safety_stock], height=bar_height, left=[forecast], color=style.PALETTE[1])
    # Label the segments inside the bar rather than in a legend: two segments and one total read
    # faster as annotations than as a key the eye has to travel to.
    for centre, amount, caption, colour in (
        (forecast / 2, forecast, "Forecast demand", "white"),
        (forecast + safety_stock / 2, safety_stock, "Safety stock (z x sigma)", "#1A1A1A"),
    ):
        axes.text(
            centre,
            0,
            f"{caption}\n{_units(amount)}",
            ha="center",
            va="center",
            fontsize=12,
            color=colour,
            fontweight="bold",
        )
    axes.text(
        target,
        0,
        f"  = {_units(target)} units",
        va="center",
        ha="left",
        fontsize=14,
        fontweight="bold",
        color="#1A1A1A",
    )
    axes.set_xlim(0, target * 1.3)
    axes.set_ylim(-0.7, 0.7)
    style.finalize(
        figure,
        title=(
            f"{product['stock_code']} {product['description']} — "
            f"{product['target_month']} recommendation"
        ),
        xlabel="Units",
        ylabel="",
        footnote=(
            f"Source: inventory_plan.csv, sigma from {product['sigma_source']} residuals, "
            f"z = {product['z']}"
        ),
    )
    return _save_deck_figure(figure, "P02_target_inventory_example")


def draw_by_abc(numbers: Numbers) -> Path:
    """Slide 8's chart: hold-out wMAPE by ABC class, champion against the main baseline."""
    models = [numbers.champion_id, numbers.baseline_id]
    subset = numbers.by_abc.loc[numbers.by_abc["model"].isin(models)]
    classes = sorted(subset["abc_class"].unique())

    style.apply_style()
    figure, axes = plt.subplots(figsize=(9.0, 4.2))
    bar_width = 0.36
    for offset, model_id in enumerate(models):
        rows = subset.loc[subset["model"] == model_id].set_index("abc_class").reindex(classes)
        positions = [index + (offset - 0.5) * bar_width for index in range(len(classes))]
        axes.bar(
            positions,
            rows["wmape"] * PERCENT_SCALE,
            width=bar_width,
            label=model_id,
            color=[style.ABC_COLORS[abc] for abc in classes],
            alpha=1.0 if offset == 0 else 0.45,
            edgecolor="#1A1A1A",
            linewidth=0.6,
        )
        for position, value in zip(positions, rows["wmape"], strict=True):
            axes.text(
                position,
                value * PERCENT_SCALE,
                f"{value * PERCENT_SCALE:.0f}",
                ha="center",
                va="bottom",
                fontsize=10,
            )
    axes.set_xticks(range(len(classes)))
    axes.set_xticklabels([f"Class {abc}" for abc in classes])
    axes.legend(title="Solid = champion, pale = baseline", frameon=False, fontsize=11)
    style.finalize(
        figure,
        title=f"Hold-out wMAPE by ABC class — {numbers.champion_id} vs {numbers.baseline_id}",
        xlabel="ABC class (training-window ABC, PRD §18.2)",
        ylabel="wMAPE (% of actual demand; lower is better)",
        footnote="Source: holdout_metrics_by_abc.csv, hold-out 2011-06 to 2011-11",
    )
    return _save_deck_figure(figure, "P03_holdout_by_abc")


def draw_inventory_kpis(numbers: Numbers) -> Path:
    """Slide 9's chart: what safety stock buys, and what it costs, for champion and baseline."""
    models = [numbers.champion_id, numbers.baseline_id]
    policies = [_POLICY_FORECAST_ONLY, _POLICY_WITH_SAFETY_STOCK]
    labels = ["Forecast only", f"Forecast + safety stock (z = {numbers.z})"]

    style.apply_style()
    figure, axes = plt.subplots(1, 2, figsize=(11.0, 4.0))
    bar_width = 0.36
    for panel, (column, title, ylabel) in enumerate(
        [
            ("fill_rate", "Fill rate — demand served", "Share of actual demand (%)"),
            ("excess_units", "Excess inventory left over", "Units"),
        ]
    ):
        for offset, (policy, label) in enumerate(zip(policies, labels, strict=True)):
            values = [numbers.kpi(model_id, policy, column) for model_id in models]
            if column == "fill_rate":
                values = [value * PERCENT_SCALE for value in values]
            positions = [index + (offset - 0.5) * bar_width for index in range(len(models))]
            axes[panel].bar(
                positions,
                values,
                width=bar_width,
                label=label,
                color=style.PALETTE[offset],
            )
        axes[panel].set_xticks(range(len(models)))
        axes[panel].set_xticklabels(models, fontsize=10)
        axes[panel].set_title(title, fontsize=13)
        axes[panel].set_ylabel(ylabel, fontsize=11)
    handles, labels_drawn = axes[0].get_legend_handles_labels()
    figure.legend(
        handles,
        labels_drawn,
        loc="upper center",
        bbox_to_anchor=(0.5, 0.98),
        ncol=len(policies),
        frameon=False,
        fontsize=10,
    )
    figure.suptitle(
        "The stockout / excess trade-off, measured in the back-test", fontsize=14, y=1.09
    )
    figure.text(
        0.5,
        -0.06,
        "Source: inventory_kpis.csv, hold-out 2011-06 to 2011-11, identical rows for every policy",
        ha="center",
        fontsize=9,
        color="#5A5A5A",
    )
    figure.tight_layout()
    return _save_deck_figure(figure, "P04_inventory_kpis")


# ---------------------------------------------------------------------------
# The deck
# ---------------------------------------------------------------------------


#: A `.pptx` is a zip archive, and python-pptx stamps each member with the wall-clock time at save.
#: Two builds of the same deck from the same run therefore differ in their bytes while every part
#: inside them is identical — which makes `git status` report a change that is not one, on a file
#: that is committed. Rewriting the archive with one timestamp taken from the *run* restores the §40
#: property the rest of the project has: same inputs, same bytes. Zip cannot store a year before
#: 1980, so a run log with an unreadable date falls back to that floor.
_ZIP_EPOCH: tuple[int, int, int, int, int, int] = (1980, 1, 1, 0, 0, 0)


def _freeze_archive_timestamps(
    path: Path, date_time: tuple[int, int, int, int, int, int]
) -> None:
    """Rewrite a saved `.pptx` so every member carries ``date_time`` instead of the clock."""
    with zipfile.ZipFile(path) as source:
        members = [(info, source.read(info.filename)) for info in source.infolist()]
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as target:
        for info, payload in members:
            frozen = zipfile.ZipInfo(info.filename, date_time=date_time)
            frozen.compress_type = info.compress_type
            frozen.external_attr = info.external_attr
            frozen.create_system = info.create_system
            target.writestr(frozen, payload)


class Deck:
    """A thin, opinionated wrapper over python-pptx: one template, one look, one message a slide."""

    def __init__(self, footnote: str, stamp: str | None) -> None:
        self.presentation = Presentation()
        self.presentation.slide_width = Emu(int(SLIDE_WIDTH))
        self.presentation.slide_height = Emu(int(SLIDE_HEIGHT))
        self._blank = self.presentation.slide_layouts[6]  # the only layout with no placeholders
        self._footnote = footnote
        self._stamp = stamp
        #: ``(title, notes)`` per slide, in order — the source of `presentation_notes.md`.
        self.notes: list[tuple[str, str]] = []

    # -- primitives ------------------------------------------------------

    def text(
        self,
        slide,  # noqa: ANN001 - a pptx Slide
        text: str,
        *,
        left: Emu,
        top: Emu,
        width: Emu,
        height: Emu,
        size: Pt,
        colour: RGBColor = INK,
        bold: bool = False,
        align: int = PP_ALIGN.LEFT,
    ):  # noqa: ANN202
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        paragraph = frame.paragraphs[0]
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = colour
        return box

    def slide(self, title: str, kicker: str | None = None):  # noqa: ANN201
        """Add a slide with the standard title, rule, kicker line and footnote."""
        slide = self.presentation.slides.add_slide(self._blank)
        self.text(
            slide,
            title,
            left=MARGIN,
            top=TITLE_TOP,
            width=CONTENT_WIDTH,
            height=TITLE_HEIGHT,
            size=TITLE_SIZE,
            bold=True,
        )
        if kicker:
            self.text(
                slide,
                kicker,
                left=MARGIN,
                top=KICKER_TOP,
                width=CONTENT_WIDTH,
                height=KICKER_HEIGHT,
                size=KICKER_SIZE,
                colour=MUTED,
            )
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN, RULE_TOP, CONTENT_WIDTH, Emu(int(RULE_HEIGHT))
        )
        rule.fill.solid()
        rule.fill.fore_color.rgb = ACCENT
        rule.line.fill.background()
        rule.shadow.inherit = False
        self.text(
            slide,
            self._footnote,
            left=MARGIN,
            top=FOOTNOTE_TOP,
            width=CONTENT_WIDTH,
            height=FOOTNOTE_HEIGHT,
            size=FOOTNOTE_SIZE,
            colour=MUTED,
        )
        if self._stamp:
            self.text(
                slide,
                self._stamp,
                left=MARGIN,
                top=FOOTNOTE_TOP - FOOTNOTE_HEIGHT,
                width=CONTENT_WIDTH,
                height=FOOTNOTE_HEIGHT,
                size=STAMP_SIZE,
                colour=WARNING,
                bold=True,
            )
        return slide

    def bullets(
        self,
        slide,  # noqa: ANN001
        items: Iterable[str],
        *,
        left: Emu = MARGIN,
        top: Emu = BODY_TOP,
        width: Emu = CONTENT_WIDTH,
        height: Emu = BODY_HEIGHT,
        size: Pt = BODY_SIZE,
        colour: RGBColor = INK,
        bullet_char: str = "• ",
    ) -> None:
        box = slide.shapes.add_textbox(left, top, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        for index, item in enumerate(items):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.space_after = Pt(9)
            run = paragraph.add_run()
            run.text = f"{bullet_char}{item}" if bullet_char else item
            run.font.size = size
            run.font.color.rgb = colour

    def heading(
        self,
        slide,  # noqa: ANN001
        text: str,
        *,
        left: Emu,
        top: Emu,
        width: Emu,
    ) -> None:
        self.text(
            slide,
            text,
            left=left,
            top=top,
            width=width,
            height=Inches(0.32),
            size=SMALL_SIZE,
            colour=ACCENT,
            bold=True,
        )

    def table(
        self,
        slide,  # noqa: ANN001
        header: Sequence[str],
        rows: Sequence[Sequence[str]],
        *,
        left: Emu = MARGIN,
        top: Emu = BODY_TOP,
        width: Emu = CONTENT_WIDTH,
        height: Emu = TABLE_HEIGHT,
        highlight: int | None = None,
    ) -> None:
        """A table with the deck's header band; ``highlight`` shades one body row (the champion)."""
        shape = slide.shapes.add_table(len(rows) + 1, len(header), left, top, width, height)
        table = shape.table
        for column, label in enumerate(header):
            cell = table.cell(0, column)
            cell.text = label
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.runs[0].font.size = TABLE_SIZE
            paragraph.runs[0].font.bold = True
            paragraph.runs[0].font.color.rgb = PAPER
        for row_index, row in enumerate(rows, start=1):
            for column, value in enumerate(row):
                cell = table.cell(row_index, column)
                cell.text = str(value)
                cell.fill.solid()
                cell.fill.fore_color.rgb = BAND if row_index - 1 == highlight else PAPER
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.runs[0].font.size = TABLE_SIZE
                paragraph.runs[0].font.bold = row_index - 1 == highlight
                paragraph.runs[0].font.color.rgb = INK

    def picture(
        self,
        slide,  # noqa: ANN001
        image: Path,
        *,
        left: Emu,
        top: Emu,
        max_width: Emu,
        max_height: Emu,
    ) -> None:
        """Place an image scaled to fit the box, centred horizontally inside it."""
        if not image.is_file():
            raise SystemExit(f"missing image for the deck: {_rel(image)}")
        pixel_width, pixel_height = Image.open(image).size
        scale = min(int(max_width) / pixel_width, int(max_height) / pixel_height)
        width, height = Emu(int(pixel_width * scale)), Emu(int(pixel_height * scale))
        offset = Emu(int((int(max_width) - int(width)) / 2))
        slide.shapes.add_picture(str(image), left + offset, top, width, height)

    def caption(self, slide, text: str, *, left: Emu, top: Emu, width: Emu) -> None:  # noqa: ANN001
        self.text(
            slide,
            text,
            left=left,
            top=top,
            width=width,
            height=Inches(0.3),
            size=SMALL_SIZE,
            colour=MUTED,
            align=PP_ALIGN.CENTER,
        )

    def note(self, slide, title: str, text: str) -> None:  # noqa: ANN001
        """Attach speaker notes to the slide and remember them for the notes file."""
        words = len(text.split())
        if words > MAX_NOTE_WORDS:
            raise SystemExit(f"speaker note for '{title}' is {words} words (max {MAX_NOTE_WORDS})")
        slide.notes_slide.notes_text_frame.text = text
        self.notes.append((title, text))

    def save(self, path: Path, stamped: tuple[int, int, int, int, int, int]) -> None:
        count = len(self.presentation.slides)
        if not SLIDE_COUNT_MIN <= count <= SLIDE_COUNT_MAX:
            raise SystemExit(
                f"the deck has {count} slides; PRD §53.2 and the course brief allow "
                f"{SLIDE_COUNT_MIN}-{SLIDE_COUNT_MAX}"
            )
        path.parent.mkdir(parents=True, exist_ok=True)
        self.presentation.save(str(path))
        _freeze_archive_timestamps(path, stamped)


# ---------------------------------------------------------------------------
# The twelve slides of PRD §53.2
# ---------------------------------------------------------------------------

#: The §53.2 outline, in order. The deck's titles start with these, and the test checks them.
SLIDE_TITLES: tuple[str, ...] = (
    "Retail Demand Forecasting & Inventory Planning",
    "The business problem",
    "Data & cleaning waterfall",
    "What happened in the data",
    "Research question & two-layer design",
    "CrewAI Flow architecture",
    "Model comparison",
    "Champion decision & results by ABC",
    "Inventory policy & back-test KPIs",
    "The product",
    "Limitations & ethics",
    "Version 2 roadmap",
)


def slide_01_title(deck: Deck, numbers: Numbers) -> None:
    run = numbers.run
    slide = deck.slide(
        SLIDE_TITLES[0],
        kicker=f"Business review · numbers as of run {run.run_id} ({run.date})",
    )
    deck.text(
        slide,
        _prd_quote(51),
        left=MARGIN,
        top=BODY_TOP,
        width=CONTENT_WIDTH,
        height=Inches(1.5),
        size=BODY_SIZE,
    )
    deck.heading(slide, "Team & roles (PRD §54)", left=MARGIN, top=Inches(3.1), width=CONTENT_WIDTH)
    deck.table(
        slide,
        ["Member", "Role"],
        [[name, role] for name, role in _team()],
        top=Inches(3.5),
        height=Inches(1.5),
    )
    provenance = run.provenance_sentence()
    deck.text(
        slide,
        f"Run status: {run.status} · mode {run.mode} · Provenance: "
        f"{provenance or 'every input table was written by this run.'}",
        left=MARGIN,
        top=Inches(5.3),
        width=CONTENT_WIDTH,
        height=Inches(1.0),
        size=SMALL_SIZE,
        colour=WARNING if provenance else MUTED,
    )
    deck.note(
        slide,
        SLIDE_TITLES[0],
        "One sentence on what we built: a Streamlit application that forecasts next month's unit "
        "demand per product and turns that forecast into a recommended target inventory. Every "
        f"number in this deck was computed by the pipeline in run {run.run_id} and read from an "
        "artifact file; nothing was typed by hand. Note the provenance line at the bottom: it "
        "says which input tables this run itself wrote and which were carried forward.",
    )


def slide_02_problem(deck: Deck, numbers: Numbers, example: Path) -> None:
    section = _prd_section(3)
    stockouts, overstock = (
        " ".join(line.split("**", maxsplit=2)[-1].strip(" —").split())
        for line in section.splitlines()
        if line.startswith("**3.")
    )
    slide = deck.slide(
        SLIDE_TITLES[1],
        kicker="Two risks, one balance — and a recommendation a planner can act on",
    )
    deck.bullets(
        slide,
        [
            f"Stockouts: {stockouts}",
            f"Overstock: {overstock}",
            "The product turns history into one actionable line per product: expected demand, the "
            "safety stock its uncertainty requires, and the target inventory level.",
            "The dataset holds no on-hand or on-order position, so the output is a Recommended "
            "Target Inventory — a level to hold, not an amount to buy (PRD §7).",
        ],
        height=Inches(2.0),
    )
    deck.picture(
        slide,
        example,
        left=MARGIN,
        top=Inches(3.5),
        max_width=CONTENT_WIDTH,
        max_height=Inches(3.1),
    )
    deck.note(
        slide,
        SLIDE_TITLES[1],
        "Thousands of products, wildly different patterns. Too little stock loses the sale; too "
        "much ties up cash and ages on the shelf. The chart is a real product from this run's "
        "inventory plan, not an illustration: the blue bar is the demand the model expects next "
        "month, the orange bar is the safety stock its forecast uncertainty requires, and the sum "
        "is the level we recommend holding. We recommend a level, never a purchase order — the "
        "data has no inventory position to subtract.",
    )


def slide_03_data(deck: Deck, numbers: Numbers) -> None:
    quality = numbers.quality
    raw = quality["raw_profile"]
    duplicates = quality["duplicates"]
    partial = quality["partial_month"]
    non_product = quality["nonproduct_codes"]
    waterfall = numbers.waterfall
    first_step, last_step = waterfall.iloc[0], waterfall.iloc[-1]
    kept_share = float(last_step["rows_after"]) / float(first_step["rows_before"])

    slide = deck.slide(
        SLIDE_TITLES[2],
        kicker="Online Retail II (UCI) — every row removed is counted, and the count is "
        "reproducible",
    )
    deck.bullets(
        slide,
        [
            f"{_units(raw['rows'])} raw sales lines, {_units(raw['distinct_stock_codes'])} stock "
            f"codes, {raw['date_range']['first'][:7]} to {raw['date_range']['last'][:7]}.",
            f"{_units(last_step['rows_after'])} lines survive cleaning ({_pct(kept_share)} of raw) "
            f"across {len(waterfall)} audited steps.",
            f"Exact duplicates are {_pct(duplicates['duplicate_row_share'], decimals=2)} of rows "
            f"and {_pct(duplicates['duplicate_units_share'], decimals=2)} of units — above the "
            "configured warning thresholds, so it is reported, not hidden.",
            f"{non_product['codes']} non-inventory codes (postage, carriage, manual adjustments) "
            "are excluded by an explicit, reviewed list.",
            f"{', '.join(partial['months'])} is partial ({_units(partial['units'])} units to the "
            f"9th): shown as a hatched actual, never scored. Last full month is "
            f"{partial['last_full_month']}.",
        ],
        top=BODY_TOP,
        height=Inches(2.2),
        width=Inches(6.0),
    )
    deck.picture(
        slide,
        paths.FIGURES_DIR / "E01_waterfall.png",
        left=MARGIN + Inches(6.2),
        top=BODY_TOP,
        max_width=CONTENT_WIDTH - Inches(6.2),
        max_height=Inches(5.0),
    )
    deck.note(
        slide,
        SLIDE_TITLES[2],
        "Two years of a UK online retailer's transactions. Cleaning is a waterfall: every step "
        "records how many rows and units it removed and why, so the count is auditable and any "
        "two runs agree. Cancellations and returns are removed as rows but never subtracted from "
        "demand — the stock still had to be on the shelf to fill the original order. December "
        "2011 stops on the 9th, so we show it hatched and never score a model on it.",
    )


def slide_04_what_happened(deck: Deck, numbers: Numbers) -> None:
    themes = ("Seasonality", "Concentration", "Lifecycle")
    sentences = _insight_sentences(_SLIDE_4_INSIGHT_TABLES, _INSIGHT_CHARS)
    slide = deck.slide(
        SLIDE_TITLES[3],
        kicker="Seasonality, concentration, lifecycle — the three facts that shaped every decision",
    )
    deck.bullets(
        slide,
        [
            f"{theme} — {sentence}"
            for theme, sentence in zip(themes, sentences, strict=True)
        ],
        top=BODY_TOP,
        height=Inches(1.7),
    )
    half = (CONTENT_WIDTH - GUTTER) / 2
    figures_top = Inches(3.45)  # clears the three full-width sentences above
    deck.picture(
        slide,
        paths.FIGURES_DIR / "E02_monthly_units.png",
        left=MARGIN,
        top=figures_top,
        max_width=Emu(int(half)),
        max_height=Inches(3.2),
    )
    deck.picture(
        slide,
        paths.FIGURES_DIR / "E06_pareto.png",
        left=MARGIN + Emu(int(half + GUTTER)),
        top=figures_top,
        max_width=Emu(int(half)),
        max_height=Inches(3.2),
    )
    deck.note(
        slide,
        SLIDE_TITLES[3],
        "Three findings drove the design. Autumn is when the money is made, so the test window "
        "covers it rather than the quiet months. Revenue is concentrated in a fifth of the "
        "catalogue, so we report accuracy by ABC class and not only in total. Most products have "
        "short lives, so the model leans on recent months rather than assuming two years of "
        "history. Each sentence comes from a computed table in the EDA report.",
    )


def slide_05_question(deck: Deck, numbers: Numbers) -> None:
    split = numbers.split
    slide = deck.slide(
        SLIDE_TITLES[4],
        kicker="One question, answered in two layers that are never mixed",
    )
    deck.text(
        slide,
        _prd_quote(1),
        left=MARGIN,
        top=BODY_TOP,
        width=CONTENT_WIDTH,
        height=Inches(1.1),
        size=Pt(17),
        bold=True,
        colour=ACCENT,
    )
    deck.bullets(
        slide,
        [
            "Layer 1 — forecasting (machine learning): one global model predicts units sold per "
            "active product for the next month. It never predicts inventory.",
            "Layer 2 — inventory policy (a deterministic rule): target inventory = "
            "ceil(max(0, forecast + z x sigma)), computed after the forecast, in code a planner "
            "can read.",
            f"Active product: at least one sale in the {numbers.active_k} months before the target "
            f"month. Forecast origin is month t-1; features may use months <= t-1 only.",
            f"Split is temporal, never shuffled: train targets {split.train_targets.start} to "
            f"{split.train_targets.end}, hold-out {split.holdout_targets.start} to "
            f"{split.holdout_targets.end}.",
            "Quarters are the sum of three one-step-ahead monthly forecasts, so quarterly error is "
            "not an independent three-month-ahead result (PRD §31-§32).",
        ],
        top=Inches(2.8),
        height=Inches(3.8),
    )
    deck.note(
        slide,
        SLIDE_TITLES[4],
        "The research question has two halves and we answer them with two different tools. How "
        "much will sell is a prediction problem, so a machine-learning model does it. How much to "
        "hold is a policy decision, so a transparent formula does it — forecast plus a safety "
        "buffer sized from the model's own past errors. Keeping them apart is what lets us change "
        "the service level without retraining anything, and it is why the model is never asked to "
        "predict inventory.",
    )


def slide_06_flow(deck: Deck, numbers: Numbers, diagram: Path) -> None:
    slide = deck.slide(
        SLIDE_TITLES[5],
        kicker="Ten steps, four validation routers, two crews of three agents — and no agent "
        "computes a number",
    )
    deck.picture(
        slide,
        diagram,
        left=MARGIN,
        top=BODY_TOP,
        max_width=Inches(7.4),
        max_height=Inches(4.9),
    )
    left = MARGIN + Inches(7.7)
    width = CONTENT_WIDTH - Inches(7.7)
    deck.heading(slide, "Data Analyst Crew (§35)", left=left, top=BODY_TOP, width=width)
    deck.bullets(
        slide,
        _crew_agents("data_analyst"),
        left=left,
        top=BODY_TOP + Inches(0.34),
        width=width,
        height=Inches(1.3),
        size=SMALL_SIZE,
    )
    deck.heading(slide, "Data Scientist Crew (§36)", left=left, top=Inches(3.2), width=width)
    deck.bullets(
        slide,
        _crew_agents("data_scientist"),
        left=left,
        top=Inches(3.5),
        width=width,
        height=Inches(1.3),
        size=SMALL_SIZE,
    )
    deck.bullets(
        slide,
        [
            "The Flow runs the deterministic tools itself; agents read the output and write prose.",
            "Every narrative is checked against the computed tables before it is accepted (§38).",
            "--no-llm and LLM runs produce numerically identical artifacts.",
        ],
        left=left,
        top=Inches(4.9),
        width=width,
        height=Inches(1.7),
        size=SMALL_SIZE,
        colour=MUTED,
    )
    deck.note(
        slide,
        SLIDE_TITLES[5],
        "The pipeline is a CrewAI Flow: ten steps from raw file to published app, with four "
        "checkpoints that can stop the run. A failed run writes a validation report and leaves the "
        "previous results untouched — it can never half-overwrite good output. Six agents in two "
        "crews review and explain, but none of them calculates anything: every number comes from "
        "code, and an agent's text is rejected if it states a number no table contains.",
    )


def slide_07_models(deck: Deck, numbers: Numbers) -> None:
    baseline = numbers.baseline_id
    improvement = numbers.improvement.set_index("model")
    rows: list[list[str]] = []
    champion_row: int | None = None
    for index, model_id in enumerate(MODEL_IDS):
        metrics = numbers.overall.loc[numbers.overall["model"] == model_id].iloc[0]
        points = improvement.loc[model_id, "wmape_points_vs_b2"]
        note = str(metrics["note"] or "")
        rows.append(
            [
                model_id + (" (champion)" if model_id == numbers.champion_id else ""),
                _pct(metrics["wmape"]),
                _pct(metrics["bias"]),
                _units(metrics["n_rows"]),
                "-" if model_id == baseline else _points(float(points)),
                "reference only, partial coverage" if note else "",
            ]
        )
        if model_id == numbers.champion_id:
            champion_row = index

    slide = deck.slide(
        SLIDE_TITLES[6],
        kicker=f"Hold-out {numbers.split.holdout_targets.start} to "
        f"{numbers.split.holdout_targets.end} — wMAPE and Bias, always together (§23)",
    )
    deck.table(
        slide,
        ["Model", "wMAPE", "Bias", "Rows scored", f"vs {baseline} (points)", "Note"],
        rows,
        top=BODY_TOP,
        height=Inches(3.5),
        highlight=champion_row,
    )
    deck.bullets(
        slide,
        [
            "wMAPE = total absolute error / total actual demand (lower is better). Bias = total "
            "(forecast - actual) / total actual: negative means we systematically under-forecast.",
            "The split is temporal and the back-test uses rolling origins, so no model ever saw "
            "the month it is scored on (§21-§22).",
            f"All models are scored on identical rows; {MODEL_IDS[2]} is reference only, because "
            "it has no month t-12 for every product.",
        ],
        top=Inches(5.0),
        height=Inches(1.7),
        size=SMALL_SIZE,
        colour=MUTED,
    )
    deck.note(
        slide,
        SLIDE_TITLES[6],
        "Read two columns together. wMAPE is the total error as a share of total demand — how far "
        "off we are. Bias is whether the misses cancel out or lean one way; a model that is "
        "accurate but always low would quietly cause stockouts. The lowest wMAPE in the table is "
        "not automatically the winner, and the next slide shows why: a strongly biased model is "
        "disqualified before accuracy is even considered.",
    )


def slide_08_champion(deck: Deck, numbers: Numbers, chart: Path) -> None:
    decision = numbers.champion
    gates = decision["config_gates"]
    failed = [
        f"{candidate['model']} (Bias {_pct(candidate['bias'])})"
        for candidate in decision["candidates"]
        if not candidate["gate1_pass"]
    ]
    improvement_points = float(decision["improvement_points"]) / PERCENT_SCALE
    verdict = (
        f"{numbers.champion_id} beats the best gate-passing baseline "
        f"({decision['best_baseline']}) by {_points(improvement_points)} wMAPE points"
    )
    verdict += (
        " — a meaningful improvement by the configured threshold."
        if decision["meaningful_improvement"]
        else " — below the configured threshold, so simple methods are competitive and we say so."
    )

    slide = deck.slide(
        SLIDE_TITLES[7],
        kicker=f"Chosen by fixed rules in code ({decision['rules_version']}), not by preference",
    )
    deck.bullets(
        slide,
        [
            f"Gate 1 — bias: |Bias| must be <= {_pct(gates['max_abs_bias'], decimals=0)}. "
            f"{len(failed)} of {len(decision['candidates'])} candidates failed: "
            f"{', '.join(failed)}.",
            f"Gate 2 — accuracy: lowest wMAPE among the survivors. Winner: "
            f"{numbers.champion_id} at {_pct(numbers.metric(numbers.champion_id, 'wmape'))}, "
            f"Bias {_pct(numbers.metric(numbers.champion_id, 'bias'))}.",
            f"Gate 4 — meaningful improvement: {verdict}",
            "Baselines are judged by the same gates. A baseline winning would be a legitimate "
            "result and would be reported as one.",
        ],
        top=BODY_TOP,
        height=Inches(2.4),
        width=Inches(6.1),
    )
    deck.picture(
        slide,
        chart,
        left=MARGIN + Inches(6.3),
        top=BODY_TOP,
        max_width=CONTENT_WIDTH - Inches(6.3),
        max_height=Inches(4.6),
    )
    deck.note(
        slide,
        SLIDE_TITLES[7],
        "The champion is picked by rules written down before the results existed. First the bias "
        "gate: any model that leans systematically high or low is out, however accurate it looks. "
        "Then the lowest error among what is left. Then a check that the winner beats the simple "
        "baseline by enough to be worth the complexity. The chart shows the same comparison split "
        "by ABC class, because the head of the catalogue is where accuracy pays.",
    )


def slide_09_inventory(deck: Deck, numbers: Numbers, chart: Path) -> None:
    champion, baseline = numbers.champion_id, numbers.baseline_id
    with_ss = _POLICY_WITH_SAFETY_STOCK
    only = _POLICY_FORECAST_ONLY
    excess_row = numbers.excess.loc[
        (numbers.excess["model"] == champion) & (numbers.excess["policy"] == with_ss)
    ].iloc[0]

    slide = deck.slide(
        SLIDE_TITLES[8],
        kicker=f"Safety stock = z x robust sigma, z = {numbers.z}; measured, not assumed",
    )
    deck.picture(
        slide,
        chart,
        left=MARGIN,
        top=BODY_TOP,
        max_width=Inches(7.0),
        max_height=Inches(4.2),
    )
    left = MARGIN + Inches(7.3)
    width = CONTENT_WIDTH - Inches(7.3)
    deck.bullets(
        slide,
        [
            f"{champion} + safety stock: fill rate "
            f"{_pct(numbers.kpi(champion, with_ss, 'fill_rate'))} (from "
            f"{_pct(numbers.kpi(champion, only, 'fill_rate'))} on the forecast alone), "
            f"{_units(numbers.kpi(champion, with_ss, 'stockout_units'))} units short, "
            f"{_units(numbers.kpi(champion, with_ss, 'excess_units'))} units excess.",
            f"{baseline} + safety stock: fill rate "
            f"{_pct(numbers.kpi(baseline, with_ss, 'fill_rate'))}, "
            f"{_units(numbers.kpi(baseline, with_ss, 'stockout_units'))} units short, "
            f"{_units(numbers.kpi(baseline, with_ss, 'excess_units'))} units excess.",
            f"Excess is concentrated: {_pct(excess_row['top_5pct_share'])} of it sits in the top "
            "5 % of product-months, so it is a short list to manage, not a spread problem.",
            f"z is a slider in the app: {', '.join(str(z) for z in numbers.z_options)}.",
            numbers.disclaimer,
        ],
        left=left,
        top=BODY_TOP,
        width=width,
        height=Inches(4.9),
        size=SMALL_SIZE,
    )
    deck.note(
        slide,
        SLIDE_TITLES[8],
        "This is the slide the business cares about. Fill rate is the share of real demand the "
        "recommended inventory would have covered; excess is what would have been left over. "
        "Adding safety stock buys a large jump in fill rate and costs excess units — the "
        "trade-off, priced. Two warnings: the multiplier does not guarantee a service level, we "
        "measure what it actually achieved; and the excess sits in a few product-months.",
    )


def slide_10_product(deck: Deck, screens: Sequence[Path]) -> None:
    slide = deck.slide(
        SLIDE_TITLES[9],
        kicker="Seven Streamlit screens; the planner never sees a model, only a recommendation",
    )
    captions = [
        "1 — Executive dashboard",
        "2 — Product forecasts & CSV export",
        "4 — Model evaluation",
        "5 — Inventory policy & z slider",
    ]
    columns = 2
    tile_width = (CONTENT_WIDTH - GUTTER) / columns
    tile_height = Inches(2.25)
    for index, (image, caption) in enumerate(zip(screens, captions, strict=True)):
        column, row = index % columns, index // columns
        left = MARGIN + Emu(int(column * (tile_width + GUTTER)))
        top = BODY_TOP + Emu(int(row * (int(tile_height) + int(Inches(0.55)))))
        deck.picture(
            slide, image, left=left, top=top, max_width=Emu(int(tile_width)), max_height=tile_height
        )
        deck.caption(slide, caption, left=left, top=top + tile_height, width=Emu(int(tile_width)))
    deck.note(
        slide,
        SLIDE_TITLES[9],
        "These are screenshots of the running application, not mock-ups. The dashboard answers "
        "'what should I hold next month' in one screen. The forecast table filters by class and "
        "exports to CSV, so a planner can work in the tool they already use. The evaluation and "
        "policy screens are the honesty layer: every metric behind the recommendation, and a "
        "slider that shows what a different service level would cost.",
    )


def slide_11_limits(deck: Deck, numbers: Numbers) -> None:
    limitations = _model_card_bullets("Limitations", _LIMITATION_TOPICS, _BULLET_CHARS)
    ethics = _model_card_bullets("Ethical considerations", _ETHICS_TOPICS, _BULLET_CHARS)
    slide = deck.slide(
        SLIDE_TITLES[10],
        kicker="Straight from the model card — the same text ships with the artifacts",
    )
    half = (CONTENT_WIDTH - GUTTER) / 2
    deck.heading(slide, "Limitations (§47)", left=MARGIN, top=BODY_TOP, width=Emu(int(half)))
    deck.bullets(
        slide,
        limitations,
        left=MARGIN,
        top=BODY_TOP + Inches(0.38),
        width=Emu(int(half)),
        height=Inches(4.6),
        size=SMALL_SIZE,
    )
    right = MARGIN + Emu(int(half + GUTTER))
    deck.heading(slide, "Ethics & licensing (§48)", left=right, top=BODY_TOP, width=Emu(int(half)))
    deck.bullets(
        slide,
        ethics,
        left=right,
        top=BODY_TOP + Inches(0.38),
        width=Emu(int(half)),
        height=Inches(4.6),
        size=SMALL_SIZE,
    )
    deck.note(
        slide,
        SLIDE_TITLES[10],
        "What this system cannot do. Twenty-four months means one prior season, so month-of-year "
        "effects rest on a single observation. There is no inventory data, so we recommend a level "
        "and not an order. New products have no history to learn from. On ethics: no decision is "
        "made about a person, customer identifiers are not used as features, and the dataset "
        "licence requires the attribution shown on every slide.",
    )


#: Slide 11 quotes the model card. §47 and §48 of the PRD name what a business audience has to be
#: told, and these are the opening words of the matching model-card bullets; the build fails if one
#: of them is no longer there.
_LIMITATION_TOPICS: tuple[str, ...] = (
    "Twenty-four months",
    "December 2011 is partial",
    "No on-hand",
    "Cold start",
    "Extreme orders",
    "Gross demand",
    "Left-censoring",
    "Service-level disclaimer",
)
_ETHICS_TOPICS: tuple[str, ...] = (
    "No decisions about individuals",
    "Customer ID is never a feature",
    "Anonymised data",
    "Residual risks",
    "Scope of the recommendation",
    "Licensing and attribution",
)
#: A slide bullet longer than this stops being a bullet and becomes a paragraph.
_BULLET_CHARS = 160
#: Slide 4's three bullets run the full slide width, so they carry more before they wrap too far.
_INSIGHT_CHARS = 210
#: The EDA tables behind slide 4's seasonality, concentration and lifecycle sentences (§35A.3).
_SLIDE_4_INSIGHT_TABLES: tuple[str, ...] = (
    "E02_sep_nov_share",
    "E06_abc_table",
    "E05_lifecycle_summary",
)
#: The roadmap has more items than fit as full-size bullets on one slide.
_ROADMAP_ITEMS = 6


def slide_12_roadmap(deck: Deck) -> None:
    slide = deck.slide(
        SLIDE_TITLES[11],
        kicker="What we would build next, in the order it would pay off (PRD §50)",
    )
    deck.bullets(
        slide,
        _roadmap_items(_ROADMAP_ITEMS),
        top=BODY_TOP,
        height=Inches(4.8),
        bullet_char="→ ",
    )
    deck.note(
        slide,
        SLIDE_TITLES[11],
        "Three of these change the answer rather than polish it. Quantile regression would give "
        "inventory levels directly instead of a forecast plus a symmetric buffer. Real inventory "
        "data would let us recommend an order amount rather than a level. Cost optimisation would "
        "replace a chosen service level with the level that actually minimises money — holding "
        "cost against lost margin. The rest is coverage: multi-horizon models, a cold-start rule, "
        "and flagging wholesale orders in evaluation.",
    )


# ---------------------------------------------------------------------------
# Assembly
# ---------------------------------------------------------------------------


def _footnote_text(run: RunProvenance) -> str:
    """The line every slide carries: which run the numbers are from, and the dataset licence."""
    return (
        f"Numbers as of run {run.run_id} ({run.date}) · Online Retail II (UCI, CC BY 4.0), "
        "Chen (2019) · generated by scripts/build_presentation.py"
    )


def _screens() -> list[Path]:
    """The four screenshots slide 10 shows, in demo order (screens 1, 2, 4, 5)."""
    wanted = [
        "S1_executive_dashboard.png",
        "S2_product_forecasts.png",
        "S4_model_evaluation.png",
        "S5_inventory_policy.png",
    ]
    missing = [name for name in wanted if not (SCREENS_DIR / name).is_file()]
    if missing:
        raise SystemExit(
            f"missing screenshots in {_rel(SCREENS_DIR)}: {', '.join(missing)} — run "
            "`python scripts/capture_screens.py` first"
        )
    return [SCREENS_DIR / name for name in wanted]


def write_notes(deck: Deck, run: RunProvenance) -> None:
    """Write `presentation_notes.md`: the speaker notes, plus what the deck was built from."""
    lines = [
        "# Speaker notes — business presentation",
        "",
        f"Generated by `scripts/build_presentation.py` from run `{run.run_id}` "
        f"({run.date}, status `{run.status}`, mode `{run.mode}`).",
        "Every number on the slides is read from a pipeline artifact; none is typed here.",
        "",
    ]
    if run.has_provenance_gap:
        lines += [
            "> **Provenance warning.** " + run.provenance_sentence(),
            ">",
            "> Tables not declared by this run: "
            + ", ".join(f"`{name}`" for name in run.undeclared)
            + ".",
            "",
        ]
    for index, (title, note) in enumerate(deck.notes, start=1):
        lines += [f"## Slide {index} — {title}", "", note, ""]
    PRESENTATION_NOTES.write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


def export_pdf() -> Path | None:
    """Convert the deck to PDF with LibreOffice, when it is installed. Optional by design."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is None:
        return None
    subprocess.run(  # noqa: S603 - a resolved executable and paths we built
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(paths.DOCS_DIR),
         str(PRESENTATION)],
        check=True,
        capture_output=True,
    )
    return PRESENTATION_PDF if PRESENTATION_PDF.is_file() else None


def build(allow_preliminary: bool = False) -> Path:
    """Read the artifacts, draw the figures, assemble the deck and write both deliverables."""
    run = RunProvenance.load()
    if not run.succeeded and not allow_preliminary:
        raise SystemExit(
            f"run {run.run_id} has status {run.status!r}, not 'success'. A deck built from it "
            "would present an earlier run's tables under this run's id. Re-run the pipeline, or "
            "pass --allow-preliminary to build a deck stamped PRELIMINARY."
        )
    numbers = Numbers.load(run)
    stamp = None if run.succeeded else f"PRELIMINARY — run {run.run_id}, status={run.status}"

    diagram = draw_flow_diagram(_flow_steps())
    example = draw_target_inventory_example(numbers)
    by_abc = draw_by_abc(numbers)
    kpi_chart = draw_inventory_kpis(numbers)

    deck = Deck(footnote=_footnote_text(run), stamp=stamp)
    slide_01_title(deck, numbers)
    slide_02_problem(deck, numbers, example)
    slide_03_data(deck, numbers)
    slide_04_what_happened(deck, numbers)
    slide_05_question(deck, numbers)
    slide_06_flow(deck, numbers, diagram)
    slide_07_models(deck, numbers)
    slide_08_champion(deck, numbers, by_abc)
    slide_09_inventory(deck, numbers, kpi_chart)
    slide_10_product(deck, _screens())
    slide_11_limits(deck, numbers)
    slide_12_roadmap(deck)

    deck.save(PRESENTATION, run.archive_timestamp)
    write_notes(deck, run)
    return PRESENTATION


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Build docs/presentation.pptx from the artifacts.")
    parser.add_argument(
        "--allow-preliminary",
        action="store_true",
        help="build even when the audited run did not succeed; stamps every slide PRELIMINARY",
    )
    parser.add_argument(
        "--pdf", action="store_true", help="also export docs/presentation.pdf (needs LibreOffice)"
    )
    args = parser.parse_args(argv)

    deck_path = build(allow_preliminary=args.allow_preliminary)
    presentation = Presentation(str(deck_path))
    print(f"wrote {_rel(deck_path)} ({len(presentation.slides)} slides)")
    print(f"wrote {_rel(PRESENTATION_NOTES)}")
    if args.pdf:
        pdf = export_pdf()
        print(f"wrote {_rel(pdf)}" if pdf else "LibreOffice not found — skipped the PDF export")
    return 0


if __name__ == "__main__":  # pragma: no cover
    raise SystemExit(main())
